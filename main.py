"""
watsonx AI Translator Agent
============================
Multi-format document translation using IBM watsonx.ai foundation models.
Exposed as an OpenAPI-compatible REST API (FastAPI).

Supported input formats: PDF, DOCX, XLSX, PPTX, HTML, Markdown, plain text
Supported language pairs: any language ↔ any language (auto-detection supported)

Usage:
  pip install -r requirements.txt
  export IBM_CLOUD_API_KEY="your-key"
  export WATSONX_PROJECT_ID="your-project-id"
  uvicorn main:app --host 0.0.0.0 --port 8000

Swagger UI:  http://localhost:8000/docs
OpenAPI JSON: http://localhost:8000/openapi.json
"""

import io
import os
import re
import tempfile
import logging
import unicodedata
from datetime import date, datetime
from enum import Enum
from pathlib import Path
from typing import Callable, Literal, Optional, Union

import boto3
import requests
from dotenv import load_dotenv
from fastapi import FastAPI, File, Form, HTTPException, Query, Request, UploadFile
from fastapi.openapi.utils import get_openapi
from fastapi.exceptions import RequestValidationError
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel, Field
from fpdf import FPDF

load_dotenv()

# ── Format-specific imports ──────────────────────────────────────────

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
from bs4 import BeautifulSoup
import pytesseract
from pdf2image import convert_from_bytes
from PIL import Image

# ── Logging ─────────────────────────────────────────────────────────

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger("watsonx-translator")
logger.info("Translator agent starting — PDF/DOCX/XLSX/PPTX/HTML/MD/TXT support active")

# ── Configuration ───────────────────────────────────────────────────

IBM_CLOUD_API_KEY    = os.getenv("IBM_CLOUD_API_KEY", "")
WATSONX_PROJECT_ID   = os.getenv("WATSONX_PROJECT_ID", "")
WATSONX_API_VERSION  = os.getenv("WATSONX_API_VERSION", "2024-05-01")
DEFAULT_WATSONX_URL  = os.getenv("WATSONX_URL", "https://us-south.ml.cloud.ibm.com")
CHUNK_SIZE           = int(os.getenv("CHUNK_SIZE", "3000"))

OUTPUT_COS_ENDPOINT   = os.getenv("OUTPUT_COS_ENDPOINT", "")
OUTPUT_COS_BUCKET     = os.getenv("OUTPUT_COS_BUCKET", "")
OUTPUT_COS_ACCESS_KEY = os.getenv("OUTPUT_COS_ACCESS_KEY") or os.getenv("AWS_ACCESS_KEY_ID", "")
OUTPUT_COS_SECRET_KEY = os.getenv("OUTPUT_COS_SECRET_KEY") or os.getenv("AWS_SECRET_ACCESS_KEY", "")

# ── Supported formats ────────────────────────────────────────────────

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".xlsx", ".xls", ".pptx", ".ppt",
    ".html", ".htm", ".md", ".txt",
}

MIME_TYPES = {
    ".pdf":  "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".xls":  "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".ppt":  "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".html": "text/html",
    ".htm":  "text/html",
    ".md":   "text/markdown",
    ".txt":  "text/plain",
}

# ── Supported Models ─────────────────────────────────────────────────

class ModelID(str, Enum):
    GRANITE_3_8B_INSTRUCT    = "ibm/granite-3-8b-instruct"
    GRANITE_3_2B_INSTRUCT    = "ibm/granite-3-2b-instruct"
    GRANITE_20B_MULTILINGUAL = "ibm/granite-20b-multilingual"
    GRANITE_13B_INSTRUCT     = "ibm/granite-13b-instruct-v2"
    LLAMA_3_1_70B_INSTRUCT   = "meta-llama/llama-3-1-70b-instruct"
    LLAMA_3_1_8B_INSTRUCT    = "meta-llama/llama-3-1-8b-instruct"
    LLAMA_3_70B_INSTRUCT     = "meta-llama/llama-3-70b-instruct"
    MISTRAL_LARGE            = "mistralai/mistral-large"
    MIXTRAL_8X7B_INSTRUCT    = "mistralai/mixtral-8x7b-instruct-v01"
    FLAN_UL2                 = "google/flan-ul2"
    ELYZA_JAPANESE_LLAMA_2   = "elyza/elyza-japanese-llama-2-7b-instruct"


class RegionURL(str, Enum):
    US_SOUTH = "https://us-south.ml.cloud.ibm.com"
    EU_DE    = "https://eu-de.ml.cloud.ibm.com"
    EU_GB    = "https://eu-gb.ml.cloud.ibm.com"
    JP_TOK   = "https://jp-tok.ml.cloud.ibm.com"


# Languages that need a CJK-capable font in the output PDF
CJK_LANGUAGES = {"japanese", "chinese", "korean", "chinese simplified", "chinese traditional"}

# Tesseract language codes for OCR (source_lang → tesseract code)
TESSERACT_LANG = {
    "japanese":   "jpn",
    "english":    "eng",
    "portuguese": "por",
    "spanish":    "spa",
    "french":     "fra",
    "german":     "deu",
    "italian":    "ita",
    "korean":     "kor",
    "chinese":    "chi_sim",
    "auto":       "jpn+eng+por+spa+fra+deu",  # try the most common packs
}


# ── CJK font path detection ──────────────────────────────────────────
# fpdf2 embeds the font directly into the PDF — no registration needed.

_CJK_FONT_CANDIDATES = [
    # TTF subfonts extracted from TTC at build time (preferred — fpdf2 renders glyphs correctly from TTF)
    "/usr/local/share/fonts/noto-cjk/NotoSansCJKjp-Regular.ttf",
    "/usr/local/share/fonts/noto-cjk/NotoSansCJKsc-Regular.ttf",
    "/usr/local/share/fonts/noto-cjk/NotoSansCJKkr-Regular.ttf",
    "/usr/local/share/fonts/noto-cjk/NotoSansCJKtc-Regular.ttf",
    # Fallbacks (TTC/OTF — glyphs may render invisible in fpdf2)
    "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
    "/usr/share/fonts/noto-cjk/NotoSansCJKjp-Regular.otf",
    "/System/Library/Fonts/PingFang.ttc",                        # macOS
    "C:/Windows/Fonts/msyh.ttc",                                 # Windows (Microsoft YaHei)
]

_cjk_font_path: Optional[str] = None

def _find_cjk_font() -> Optional[str]:
    global _cjk_font_path
    if _cjk_font_path:
        return _cjk_font_path
    import glob
    paths = _CJK_FONT_CANDIDATES + \
            glob.glob("/usr/local/share/fonts/noto-cjk/*.ttf") + \
            glob.glob("/usr/share/fonts/**/*CJK*.otf", recursive=True)
    for p in paths:
        if os.path.exists(p):
            _cjk_font_path = p
            logger.info("CJK font found: %s", p)
            return p
    logger.warning("No CJK font found — install fonts-noto-cjk for Japanese/Chinese/Korean PDF output")
    return None


# ── Dynamic prompt builder ───────────────────────────────────────────

def _build_prompt(text: str, source_lang: str, target_lang: str, model_id: str) -> str:
    """Build a model-family–aware translation prompt for any language pair."""
    if source_lang.strip().lower() == "auto":
        instruction = (
            f"Detect the source language and translate the following text to {target_lang}. "
            "Preserve paragraph structure and all formatting. "
            "Output ONLY the translation, with no explanations or notes."
        )
    else:
        instruction = (
            f"Translate the following {source_lang} text to {target_lang}. "
            "Preserve paragraph structure and all formatting. "
            "Output ONLY the translation, with no explanations or notes."
        )

    model_lower = model_id.lower()
    if "granite" in model_lower:
        return f"<|system|>\n{instruction}\n<|user|>\n{text}\n<|assistant|>\n"
    elif "llama" in model_lower:
        return (
            f"<|begin_of_text|><|start_header_id|>system<|end_header_id|>\n\n"
            f"{instruction}<|eot_id|><|start_header_id|>user<|end_header_id|>\n\n"
            f"{text}<|eot_id|><|start_header_id|>assistant<|end_header_id|>\n\n"
        )
    elif "mistral" in model_lower or "mixtral" in model_lower:
        return f"[INST] {instruction}\n\n{text} [/INST]"
    return f"{instruction}\n\nText:\n{text}\n\nTranslation:\n"


# ── Response / Request schemas ───────────────────────────────────────

class ModelInfo(BaseModel):
    id: str
    name: str
    family: str


class ModelsResponse(BaseModel):
    models: list[ModelInfo]


class FormatInfo(BaseModel):
    extension: str
    mime_type: str
    requires: str


class FormatsResponse(BaseModel):
    formats: list[FormatInfo]


class TranslationPageDetail(BaseModel):
    page: int
    source_chars: int
    translated_chars: int


class TranslateResponse(BaseModel):
    message: str
    source_lang: str
    target_lang: str
    pages_translated: int
    model_used: str
    region: str
    download_url: str
    pages: list[TranslationPageDetail]


class TranslateTextRequest(BaseModel):
    text: str = Field(..., description="Text to translate")
    source_lang: Optional[str] = Field(default="auto", description="Source language (e.g. 'Japanese'). Use 'auto' to detect.")
    target_lang: Optional[str] = Field(default="English", description="Target language (e.g. 'English', 'Spanish')")
    model_id: Optional[str] = Field(default=None, description="watsonx model ID")
    region: Optional[str] = Field(default=None, description="watsonx region URL")
    filename: Optional[str] = Field(default=None, description="Output filename stem")


class TranslateTextResponse(BaseModel):
    translated_text: str
    source_lang: str
    target_lang: str
    model_used: str
    source_chars: int
    translated_chars: int
    download_url: Optional[str] = None


class HealthResponse(BaseModel):
    status: str
    watsonx_url: str
    project_configured: bool
    formats_available: list[str]
    cjk_font_path: Optional[str]     # None = CJK PDF output will be blank
    pdf_engine: str                   # "fpdf2" or "reportlab"


# ── Upload-source schemas (backward compat) ──────────────────────────

class FilePathSource(BaseModel):
    type: Literal["file_path"]
    path: str = Field(..., description="Absolute or relative path to the document on the server filesystem")


class URLSource(BaseModel):
    type: Literal["url"]
    url: str = Field(..., description="Publicly accessible or pre-signed URL pointing to the document")
    headers: Optional[dict] = Field(default=None, description="Optional HTTP headers (e.g. Authorization)")


class BucketSource(BaseModel):
    type: Literal["bucket"]
    endpoint_url: Optional[str] = Field(default=None)
    bucket: str
    key: str
    access_key_id: Optional[str] = Field(default=None)
    secret_access_key: Optional[str] = Field(default=None)
    region_name: Optional[str] = Field(default=None)


class TranslateDocumentBase64Request(BaseModel):
    """
    Request body for AI orchestration tools (e.g. watsonx Orchestrate).
    The orchestration platform base64-encodes the user's file attachment automatically
    before calling this skill — the end user just uploads a normal file.
    """
    file: str = Field(
        ...,
        description=(
            "Base64-encoded content of the document to translate. "
            "Supported formats: PDF, DOCX, XLSX, PPTX, HTML, Markdown, TXT. "
            "AI orchestration platforms (e.g. watsonx Orchestrate) encode the user's "
            "file attachment automatically — the end user just attaches a normal file."
        ),
        json_schema_extra={"format": "byte"},
    )
    filename: Optional[str] = Field(
        default="document.pdf",
        description=(
            "Original filename including extension (e.g. 'report.docx', 'invoice.pdf'). "
            "Used to detect the file format and name the output file. "
            "Always include the correct extension."
        ),
    )
    source_lang: Optional[str] = Field(
        default="auto",
        description=(
            "Language of the source document. Examples: 'Japanese', 'Portuguese', 'Spanish', 'French'. "
            "Use 'auto' to let the model detect the language automatically."
        ),
    )
    target_lang: Optional[str] = Field(
        default="English",
        description="Language to translate into. Examples: 'English', 'Japanese', 'Spanish', 'French'.",
    )
    model_id: Optional[str] = Field(
        default=None,
        description="watsonx.ai model ID. Defaults to ibm/granite-3-8b-instruct.",
    )
    region: Optional[str] = Field(default=None, description="watsonx.ai region URL.")
    project_id: Optional[str] = Field(default=None, description="watsonx project ID.")


# Backward-compat alias used by existing integrations
TranslatePdfBase64Request = TranslateDocumentBase64Request


class TranslateFromSourceRequest(BaseModel):
    source: Union[FilePathSource, URLSource, BucketSource] = Field(..., discriminator="type")
    source_lang: Optional[str] = Field(default="Japanese", description="Source language")
    target_lang: Optional[str] = Field(default="English", description="Target language")
    model_id: Optional[str] = Field(default=None)
    region: Optional[str] = Field(default=None)
    project_id: Optional[str] = Field(default=None)


# ── IAM Token Manager ────────────────────────────────────────────────

class IAMTokenManager:
    def __init__(self):
        self._token: Optional[str] = None

    def get_token(self, api_key: Optional[str] = None) -> str:
        key = api_key or IBM_CLOUD_API_KEY
        if not key:
            raise HTTPException(
                status_code=500,
                detail="IBM_CLOUD_API_KEY not configured.",
            )
        logger.info("Requesting IAM token...")
        resp = requests.post(
            "https://iam.cloud.ibm.com/identity/token",
            headers={"Content-Type": "application/x-www-form-urlencoded"},
            data=f"grant_type=urn:ibm:params:oauth:grant-type:apikey&apikey={key}",
            timeout=30,
        )
        if resp.status_code != 200:
            raise HTTPException(status_code=502, detail=f"IAM token request failed: {resp.text}")
        self._token = resp.json()["access_token"]
        logger.info("IAM token acquired.")
        return self._token


token_manager = IAMTokenManager()


# ── Document extraction ──────────────────────────────────────────────

def _ocr_pdf(file_bytes: bytes, source_lang: str = "auto") -> list[str]:
    """OCR fallback: convert each PDF page to an image and run Tesseract."""
    lang_code = TESSERACT_LANG.get(source_lang.lower(), "eng")
    logger.info("OCR fallback: converting PDF pages to images (lang=%s)...", lang_code)
    try:
        images = convert_from_bytes(file_bytes, dpi=300)
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Could not rasterize PDF for OCR: {exc}")

    pages = []
    for i, img in enumerate(images):
        text = pytesseract.image_to_string(img, lang=lang_code)
        if text and text.strip():
            pages.append(text.strip())
            logger.info("OCR page %d: %d chars", i + 1, len(text))
        else:
            logger.warning("OCR page %d: no text detected", i + 1)
    return pages


def _extract_pdf(file_bytes: bytes, source_lang: str = "auto") -> list[str]:
    """Extract text from a PDF. Tries the text layer first; falls back to OCR for scanned pages."""
    try:
        reader = PdfReader(io.BytesIO(file_bytes))
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Failed to read PDF: {exc}")

    pages = []
    scanned_pages: list[int] = []          # page indices with no text layer

    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            pages.append((i, text.strip()))
            logger.info("Page %d: %d chars (text layer)", i + 1, len(text))
        else:
            pages.append((i, None))
            scanned_pages.append(i)
            logger.warning("Page %d: no text layer — will OCR", i + 1)

    # OCR only the pages that had no text
    if scanned_pages:
        logger.info("Running OCR on %d scanned page(s)...", len(scanned_pages))
        try:
            images = convert_from_bytes(file_bytes, dpi=300)
            lang_code = TESSERACT_LANG.get(source_lang.lower(), "eng")
            for idx in scanned_pages:
                if idx < len(images):
                    ocr_text = pytesseract.image_to_string(images[idx], lang=lang_code)
                    if ocr_text and ocr_text.strip():
                        pages[idx] = (idx, ocr_text.strip())
                        logger.info("OCR page %d: %d chars", idx + 1, len(ocr_text))
        except Exception as exc:
            logger.error("OCR failed: %s — scanned pages will be skipped", exc)

    result = [text for _, text in sorted(pages) if text]
    return result


def _extract_html(file_bytes: bytes) -> list[str]:
    """Extract visible text from an HTML file."""
    soup = BeautifulSoup(file_bytes.decode("utf-8", errors="replace"), "html.parser")
    text = soup.get_text(separator="\n").strip()
    return [text] if text else []


def _extract_text(file_bytes: bytes) -> list[str]:
    """Extract plain text (Markdown, TXT)."""
    text = file_bytes.decode("utf-8", errors="replace").strip()
    return [text] if text else []


def extract_pages(file_bytes: bytes, filename: str, source_lang: str = "auto") -> list[str]:
    """Extract text from a document for the PDF-rebuild path (PDF, HTML, MD, TXT)."""
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(file_bytes, source_lang)
    if ext in (".html", ".htm"):
        return _extract_html(file_bytes)
    if ext in (".md", ".txt"):
        return _extract_text(file_bytes)
    raise HTTPException(status_code=415, detail=f"Unsupported format for text extraction: '{ext}'")


# ── Text chunking ────────────────────────────────────────────────────

def chunk_text(text: str, max_chars: int = CHUNK_SIZE) -> list[str]:
    if len(text) <= max_chars:
        return [text]

    chunks: list[str] = []
    current = ""
    for para in text.split("\n"):
        if len(current) + len(para) + 1 > max_chars and current:
            chunks.append(current.strip())
            current = para
        else:
            current += ("\n" if current else "") + para

    if current.strip():
        chunks.append(current.strip())

    final: list[str] = []
    for chunk in chunks:
        if len(chunk) > max_chars:
            for i in range(0, len(chunk), max_chars):
                final.append(chunk[i: i + max_chars])
        else:
            final.append(chunk)
    return final


# ── Translation core ─────────────────────────────────────────────────

def translate_text(
    text: str,
    model_id: str,
    token: str,
    watsonx_url: str,
    project_id: str,
    source_lang: str = "auto",
    target_lang: str = "English",
    temperature: float = 0.1,
    max_new_tokens: int = 4096,
) -> str:
    prompt = _build_prompt(text, source_lang, target_lang, model_id)
    payload = {
        "model_id": model_id,
        "input": prompt,
        "project_id": project_id,
        "parameters": {
            "max_new_tokens": max_new_tokens,
            "temperature": temperature,
            "top_p": 0.95,
            "top_k": 50,
            "repetition_penalty": 1.05,
            "stop_sequences": [],
        },
    }
    url = f"{watsonx_url}/ml/v1/text/generation?version={WATSONX_API_VERSION}"
    logger.info("watsonx call: model=%s chars=%d %s→%s", model_id, len(text), source_lang, target_lang)

    resp = requests.post(
        url,
        headers={
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/json",
            "Accept": "application/json",
        },
        json=payload,
        timeout=180,
    )
    if resp.status_code != 200:
        logger.error("watsonx error %d: %s", resp.status_code, resp.text)
        raise HTTPException(status_code=502, detail=f"watsonx.ai API error ({resp.status_code}): {resp.text[:500]}")

    results = resp.json().get("results", [])
    if not results:
        raise HTTPException(status_code=502, detail="Empty response from watsonx.ai")

    generated = results[0].get("generated_text", "").strip()
    logger.info("Translation received: %d chars", len(generated))
    return generated


def translate_page(
    page_text: str,
    model_id: str,
    token: str,
    watsonx_url: str,
    project_id: str,
    source_lang: str = "auto",
    target_lang: str = "English",
) -> str:
    """Translate a full page, chunking long content automatically."""
    chunks = chunk_text(page_text)
    translated = []
    for i, chunk in enumerate(chunks):
        logger.info("  chunk %d/%d (%d chars)", i + 1, len(chunks), len(chunk))
        translated.append(translate_text(chunk, model_id, token, watsonx_url, project_id, source_lang, target_lang))
    return "\n\n".join(translated)


# ── In-place format reconstruction ───────────────────────────────────

def _translate_fn(
    model_id: str, token: str, watsonx_url: str, project_id: str,
    source_lang: str, target_lang: str,
) -> Callable[[str], str]:
    """Return a callable that translates a string segment."""
    def _fn(text: str) -> str:
        if not text or not text.strip():
            return text
        return translate_text(text, model_id, token, watsonx_url, project_id, source_lang, target_lang)
    return _fn


def _replace_para_text(para, translated: str) -> None:
    """Replace all runs in a python-docx paragraph with a single translated run."""
    if not para.runs:
        para.add_run(translated)
        return
    para.runs[0].text = translated
    for run in para.runs[1:]:
        run.text = ""


def translate_docx(file_bytes: bytes, fn: Callable[[str], str]) -> bytes:
    if not DOCX_AVAILABLE:
        raise HTTPException(status_code=415, detail="python-docx not installed. pip install python-docx")

    doc = DocxDocument(io.BytesIO(file_bytes))

    for para in doc.paragraphs:
        if para.text.strip():
            _replace_para_text(para, fn(para.text))

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    if para.text.strip():
                        _replace_para_text(para, fn(para.text))

    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


def translate_pptx(file_bytes: bytes, fn: Callable[[str], str]) -> bytes:
    if not PPTX_AVAILABLE:
        raise HTTPException(status_code=415, detail="python-pptx not installed. pip install python-pptx")

    prs = Presentation(io.BytesIO(file_bytes))

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                full = "".join(r.text for r in para.runs)
                if not full.strip():
                    continue
                translated = fn(full)
                if para.runs:
                    para.runs[0].text = translated
                    for run in para.runs[1:]:
                        run.text = ""
                else:
                    para.add_run(translated)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def translate_xlsx(file_bytes: bytes, fn: Callable[[str], str]) -> bytes:
    if not XLSX_AVAILABLE:
        raise HTTPException(status_code=415, detail="openpyxl not installed. pip install openpyxl")

    wb = load_workbook(io.BytesIO(file_bytes))

    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and cell.value.strip():
                    cell.value = fn(cell.value)

    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


# ── PDF generation (ReportLab) ───────────────────────────────────────

def build_translated_pdf(
    pages: list[str],
    output_path: str,
    source_lang: str = "auto",
    target_lang: str = "English",
) -> None:
    """Build a translated PDF using fpdf2 with proper Unicode/CJK font embedding."""
    needs_cjk = target_lang.strip().lower() in CJK_LANGUAGES
    cjk_path  = _find_cjk_font() if needs_cjk else None

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=20)
    pdf.set_margins(20, 20, 20)

    _DEJAVU = "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
    if cjk_path:
        pdf.add_font("body", fname=cjk_path)
        body_font = "body"
    elif os.path.exists(_DEJAVU):
        pdf.add_font("body", fname=_DEJAVU)
        body_font = "body"
    else:
        body_font = "Helvetica"
        if needs_cjk:
            logger.warning("CJK font unavailable — output may show boxes for %s", target_lang)

    _header_font = "body" if body_font == "body" else "Helvetica"
    src = source_lang if source_lang.lower() != "auto" else "auto-detected"

    for i, page_text in enumerate(pages):
        pdf.add_page()

        # Page header
        pdf.set_font(_header_font, size=9)
        pdf.set_text_color(120, 120, 120)
        pdf.cell(
            0, 8,
            f"Translated Document  |  {src} -> {target_lang}  |  Page {i + 1} of {len(pages)}",
            new_x="LMARGIN", new_y="NEXT",
        )
        pdf.ln(4)
        pdf.set_draw_color(200, 200, 200)
        pdf.line(pdf.l_margin, pdf.get_y(), pdf.w - pdf.r_margin, pdf.get_y())
        pdf.ln(6)

        # Body text
        pdf.set_font(body_font, size=11)
        pdf.set_text_color(0, 0, 0)
        # Strip control characters but keep all Unicode (CJK, Arabic, etc.)
        clean = "".join(
            c for c in page_text
            if c in ("\n", "\t", " ") or not unicodedata.category(c).startswith("C")
        )
        pdf.multi_cell(0, 7, clean)

    pdf.output(output_path)
    logger.info("PDF built: %s (font=%s)", output_path, body_font)


# ── COS / local file management ──────────────────────────────────────

def upload_to_cos(local_path: str, object_key: str, mime_type: str = "application/pdf") -> str:
    s3 = boto3.client(
        "s3",
        endpoint_url=OUTPUT_COS_ENDPOINT,
        aws_access_key_id=OUTPUT_COS_ACCESS_KEY,
        aws_secret_access_key=OUTPUT_COS_SECRET_KEY,
    )
    with open(local_path, "rb") as fh:
        s3.upload_fileobj(
            fh, OUTPUT_COS_BUCKET, object_key,
            ExtraArgs={"ContentType": mime_type, "ACL": "public-read"},
        )
    logger.info("Uploaded to COS: %s/%s", OUTPUT_COS_BUCKET, object_key)
    host = OUTPUT_COS_ENDPOINT.replace("https://", "").replace("http://", "")
    return f"https://{OUTPUT_COS_BUCKET}.{host}/{object_key}"


def _output_filename(original_name: str, output_ext: Optional[str] = None) -> str:
    stem = re.sub(r"[^a-zA-Z0-9_\-]", "_", Path(original_name).stem) if original_name else "document"
    ext = output_ext or Path(original_name).suffix.lower() or ".pdf"
    now = datetime.now().strftime("%Y%m%d_%H%M%S")
    return f"{stem}_translated_{now}{ext}"


def save_and_finalize(
    content: bytes,
    original_name: str,
    output_ext: str,
    base_url: str = "",
) -> str:
    """Write bytes to temp file, upload to COS (if configured) or serve locally."""
    filename = _output_filename(original_name, output_ext)
    path = os.path.join(tempfile.gettempdir(), filename)
    with open(path, "wb") as f:
        f.write(content)

    if OUTPUT_COS_ENDPOINT and OUTPUT_COS_BUCKET and OUTPUT_COS_ACCESS_KEY:
        mime = MIME_TYPES.get(output_ext, "application/octet-stream")
        return upload_to_cos(path, f"translated/{filename}", mime)

    prefix = base_url.rstrip("/") if base_url else ""
    return f"{prefix}/api/v1/download/{filename}"


def load_bytes_from_source(source: Union[FilePathSource, URLSource, BucketSource]) -> tuple[bytes, str]:
    """Fetch document bytes and original filename from any source type."""
    if source.type == "file_path":
        path = source.path
        if not os.path.exists(path):
            raise HTTPException(status_code=404, detail=f"File not found: {path}")
        with open(path, "rb") as f:
            return f.read(), os.path.basename(path)

    if source.type == "url":
        headers = source.headers or {}
        resp = requests.get(source.url, headers=headers, timeout=60)
        if resp.status_code != 200:
            raise HTTPException(status_code=502, detail=f"Failed to download URL ({resp.status_code})")
        filename = source.url.split("?")[0].rstrip("/").split("/")[-1] or "document"
        return resp.content, filename

    # bucket
    s3 = boto3.client(
        "s3",
        endpoint_url=source.endpoint_url,
        aws_access_key_id=source.access_key_id or OUTPUT_COS_ACCESS_KEY,
        aws_secret_access_key=source.secret_access_key or OUTPUT_COS_SECRET_KEY,
        region_name=source.region_name,
    )
    buf = io.BytesIO()
    try:
        s3.download_fileobj(source.bucket, source.key, buf)
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"Bucket download failed: {exc}")
    return buf.getvalue(), source.key.split("/")[-1]


# ── FastAPI app ──────────────────────────────────────────────────────

app = FastAPI(
    title="watsonx AI Translator",
    description=(
        "Multi-format document translation using IBM watsonx.ai.\n\n"
        "**Supported formats:** PDF, DOCX, XLSX, PPTX, HTML, Markdown, plain text\n\n"
        "**Language pairs:** Any ↔ Any (Japanese, English, Portuguese, Spanish, French, German, …)"
    ),
    version="2.0.0",
    docs_url="/docs",
    redoc_url="/redoc",
)


def _custom_openapi():
    if app.openapi_schema:
        return app.openapi_schema
    schema = get_openapi(
        title=app.title,
        version=app.version,
        openapi_version="3.0.3",
        description=app.description,
        routes=app.routes,
    )
    app_url = os.getenv("APP_URL", "").rstrip("/")
    if app_url:
        schema["servers"] = [{"url": app_url}]
    app.openapi_schema = schema
    return app.openapi_schema

app.openapi = _custom_openapi


@app.exception_handler(RequestValidationError)
async def validation_error_handler(request: Request, exc: RequestValidationError):
    logger.error("422 Validation error on %s %s: %s", request.method, request.url, exc.errors())
    return JSONResponse(status_code=422, content={"detail": exc.errors()})


@app.exception_handler(Exception)
async def global_exception_handler(request: Request, exc: Exception):
    if isinstance(exc, HTTPException):
        return JSONResponse(status_code=exc.status_code, content={"detail": exc.detail})
    logger.exception("Unhandled error: %s", exc)
    return JSONResponse(status_code=500, content={"detail": "Internal server error"})


# ── Info endpoints ───────────────────────────────────────────────────

@app.get("/api/v1/health", response_model=HealthResponse, tags=["Info"])
async def health_check():
    return HealthResponse(
        status="ok",
        watsonx_url=DEFAULT_WATSONX_URL,
        project_configured=bool(WATSONX_PROJECT_ID),
        formats_available=sorted(SUPPORTED_EXTENSIONS),
        cjk_font_path=_find_cjk_font(),
        pdf_engine="fpdf2",
    )


@app.get("/api/v1/models", response_model=ModelsResponse, tags=["Info"])
async def list_models():
    models = []
    for m in ModelID:
        parts = m.value.split("/")
        family = parts[0] if len(parts) > 1 else "unknown"
        name   = parts[1] if len(parts) > 1 else m.value
        models.append(ModelInfo(id=m.value, name=name, family=family))
    return ModelsResponse(models=models)


@app.get("/api/v1/formats", response_model=FormatsResponse, tags=["Info"])
async def list_formats():
    lib_map = {
        ".pdf":  "docling (OCR-capable) or pypdf",
        ".docx": "python-docx",
        ".xlsx": "openpyxl",
        ".xls":  "openpyxl",
        ".pptx": "python-pptx",
        ".ppt":  "python-pptx (converted to pptx)",
        ".html": "docling",
        ".htm":  "docling",
        ".md":   "docling",
        ".txt":  "docling",
    }
    available = []
    for ext, mime in MIME_TYPES.items():
        available.append(FormatInfo(extension=ext, mime_type=mime, requires=lib_map.get(ext, "docling")))
    return FormatsResponse(formats=available)


@app.get("/api/v1/regions", tags=["Info"])
async def list_regions():
    return {"regions": [{"id": r.name.lower(), "url": r.value} for r in RegionURL]}


# ── Unified translation endpoint ─────────────────────────────────────

@app.post(
    "/api/v1/translate/document",
    response_model=TranslateResponse,
    tags=["Translation"],
    summary="Translate any document (PDF, Word, Excel, PowerPoint, …)",
    description=(
        "Upload a document in any supported format. The agent extracts text, translates it "
        "using the selected watsonx.ai model, and returns the translated file in the **same format** "
        "as the input (DOCX→DOCX, XLSX→XLSX, PPTX→PPTX, PDF→PDF).\n\n"
        "Set `source_lang='auto'` for automatic language detection.\n\n"
        "**Examples:** Japanese→English, English→Japanese, Portuguese→Spanish, French→German, …"
    ),
)
async def translate_document(
    request: Request,
    file: UploadFile = File(..., description="Document to translate"),
    source_lang: str = Form(default="auto", description="Source language (e.g. 'Japanese', 'Portuguese'). Use 'auto' to detect automatically."),
    target_lang: str = Form(default="English", description="Target language (e.g. 'English', 'Spanish', 'French', 'Japanese')"),
    model_id: str = Form(default=ModelID.GRANITE_3_8B_INSTRUCT.value, description="watsonx.ai model ID"),
    region: Optional[str] = Form(default=None, description="watsonx.ai region URL"),
    project_id: Optional[str] = Form(default=None, description="watsonx project ID"),
):
    filename = file.filename or "document"
    ext = Path(filename).suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        raise HTTPException(
            status_code=415,
            detail=f"Unsupported format '{ext}'. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}",
        )

    watsonx_url  = region or DEFAULT_WATSONX_URL
    wx_project_id = project_id or WATSONX_PROJECT_ID
    if not wx_project_id:
        raise HTTPException(status_code=400, detail="Project ID required. Set WATSONX_PROJECT_ID env var or pass as query param.")

    file_bytes = await file.read()
    logger.info("Received document: %s (%d bytes) | %s → %s", filename, len(file_bytes), source_lang, target_lang)

    token = token_manager.get_token()
    fn = _translate_fn(model_id, token, watsonx_url, wx_project_id, source_lang, target_lang)

    page_details: list[TranslationPageDetail] = []
    output_ext = ext

    # ── Format dispatch ──────────────────────────────────────────────

    if ext == ".docx":
        output_bytes = translate_docx(file_bytes, fn)

    elif ext in (".xlsx", ".xls"):
        output_bytes = translate_xlsx(file_bytes, fn)
        output_ext = ".xlsx"

    elif ext in (".pptx", ".ppt"):
        output_bytes = translate_pptx(file_bytes, fn)
        output_ext = ".pptx"

    else:
        # PDF and text-based formats: extract → translate pages → rebuild as PDF
        pages = extract_pages(file_bytes, filename, source_lang)
        if not pages:
            raise HTTPException(
                status_code=422,
                detail=(
                    "No extractable text found in this document. "
                    "For scanned PDFs, ensure OCR has been applied first "
                    "(e.g. with 'ocrmypdf'). "
                    "For image-only PDFs, docling OCR must be enabled."
                ),
            )

        translated_pages: list[str] = []
        for i, page_text in enumerate(pages):
            logger.info("Translating page %d/%d...", i + 1, len(pages))
            translated = translate_page(page_text, model_id, token, watsonx_url, wx_project_id, source_lang, target_lang)
            translated_pages.append(translated)
            page_details.append(TranslationPageDetail(
                page=i + 1,
                source_chars=len(page_text),
                translated_chars=len(translated),
            ))

        # Rebuild as PDF regardless of input format (HTML, MD, TXT → PDF)
        output_ext = ".pdf"
        tmp_pdf = os.path.join(tempfile.gettempdir(), _output_filename(filename, ".pdf"))
        build_translated_pdf(translated_pages, tmp_pdf, source_lang, target_lang)
        with open(tmp_pdf, "rb") as f:
            output_bytes = f.read()

    download_url = save_and_finalize(output_bytes, filename, output_ext, str(request.base_url))
    logger.info("Done: %d bytes → %s", len(output_bytes), download_url)

    return TranslateResponse(
        message=f"Translation complete ({source_lang} → {target_lang})",
        source_lang=source_lang,
        target_lang=target_lang,
        pages_translated=len(page_details) or 1,
        model_used=model_id,
        region=watsonx_url,
        download_url=download_url,
        pages=page_details,
    )


# ── Legacy PDF endpoints (backward compatible, now language-aware) ────

@app.post(
    "/api/v1/translate/pdf",
    response_model=TranslateResponse,
    tags=["Translation"],
    summary="Translate a PDF document",
    description=(
        "Upload a PDF. Extracts text (using docling with pypdf fallback), translates each page, "
        "and returns a downloadable translated PDF.\n\n"
        "Now supports any language pair — not just Japanese→English."
    ),
)
async def translate_pdf(
    request: Request,
    file: UploadFile = File(..., description="PDF file to translate"),
    source_lang: str = Query(default="Japanese", description="Source language"),
    target_lang: str = Query(default="English", description="Target language"),
    model_id: str = Query(default=ModelID.GRANITE_3_8B_INSTRUCT.value),
    region: Optional[str] = Query(default=None),
    project_id: Optional[str] = Query(default=None),
):
    filename = file.filename or "document.pdf"
    if not filename.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are accepted by this endpoint. Use /api/v1/translate/document for other formats.")

    watsonx_url   = region or DEFAULT_WATSONX_URL
    wx_project_id = project_id or WATSONX_PROJECT_ID
    if not wx_project_id:
        raise HTTPException(status_code=400, detail="Project ID required.")

    file_bytes = await file.read()
    logger.info("PDF upload: %s (%d bytes) | %s → %s", filename, len(file_bytes), source_lang, target_lang)

    pages = extract_pages(file_bytes, filename, source_lang)
    if not pages:
        raise HTTPException(
            status_code=422,
            detail="No extractable text found in the PDF. The file may be a scanned image. Use ocrmypdf to add a text layer first.",
        )
    logger.info("Extracted %d pages.", len(pages))

    token = token_manager.get_token()
    translated_pages: list[str] = []
    page_details: list[TranslationPageDetail] = []

    for i, page_text in enumerate(pages):
        logger.info("Translating page %d/%d...", i + 1, len(pages))
        translated = translate_page(page_text, model_id, token, watsonx_url, wx_project_id, source_lang, target_lang)
        translated_pages.append(translated)
        page_details.append(TranslationPageDetail(page=i + 1, source_chars=len(page_text), translated_chars=len(translated)))

    tmp_pdf = os.path.join(tempfile.gettempdir(), _output_filename(filename, ".pdf"))
    build_translated_pdf(translated_pages, tmp_pdf, source_lang, target_lang)
    with open(tmp_pdf, "rb") as f:
        pdf_bytes = f.read()

    download_url = save_and_finalize(pdf_bytes, filename, ".pdf", str(request.base_url))
    logger.info("Translation complete: %d pages → %s", len(translated_pages), download_url)

    return TranslateResponse(
        message=f"Translation complete ({source_lang} → {target_lang})",
        source_lang=source_lang,
        target_lang=target_lang,
        pages_translated=len(translated_pages),
        model_used=model_id,
        region=watsonx_url,
        download_url=download_url,
        pages=page_details,
    )


@app.post(
    "/api/v1/translate/document-base64",
    response_model=TranslateResponse,
    tags=["Translation"],
    summary="Translate any document — base64 JSON body (for AI orchestration)",
    description=(
        "**Designed for AI orchestration platforms such as watsonx Orchestrate.**\n\n"
        "The end user simply attaches a file in the chat. "
        "The orchestration platform automatically base64-encodes the attachment and calls this skill "
        "— the user never deals with base64 themselves.\n\n"
        "Supported formats: **PDF, DOCX, XLSX, PPTX, HTML, Markdown, TXT**\n\n"
        "DOCX → DOCX, XLSX → XLSX, PPTX → PPTX (format preserved). PDF/HTML/MD/TXT → PDF.\n\n"
        "Set `source_lang='auto'` for automatic language detection.\n\n"
        "**Example payload:**\n"
        "```json\n"
        "{\n"
        '  "file": "<base64-encoded content>",\n'
        '  "filename": "report.docx",\n'
        '  "source_lang": "Portuguese",\n'
        '  "target_lang": "Japanese"\n'
        "}\n"
        "```"
    ),
)
async def translate_document_base64(request: Request, body: TranslateDocumentBase64Request):
    return await _translate_base64_impl(request, body)


@app.post(
    "/api/v1/translate/pdf-base64",
    response_model=TranslateResponse,
    tags=["Translation"],
    summary="[Deprecated] Use /api/v1/translate/document-base64 instead",
    description="Backward-compatible alias. Use `/api/v1/translate/document-base64` for new integrations.",
    include_in_schema=True,
)
async def translate_pdf_base64(request: Request, body: TranslatePdfBase64Request):
    return await _translate_base64_impl(request, body)


async def _translate_base64_impl(request: Request, body: TranslateDocumentBase64Request):
    import base64

    model_id      = body.model_id or ModelID.GRANITE_3_8B_INSTRUCT.value
    watsonx_url   = body.region or DEFAULT_WATSONX_URL
    wx_project_id = body.project_id or WATSONX_PROJECT_ID
    source_lang   = body.source_lang or "Japanese"
    target_lang   = body.target_lang or "English"

    if not wx_project_id:
        raise HTTPException(status_code=400, detail="Project ID required.")

    try:
        raw = body.file
        logger.info("Base64 input: len=%d prefix=%r", len(raw), raw[:80])
        # Strip data-URL prefix: "data:<mime>;base64,<data>"
        if "," in raw and raw.startswith("data:"):
            raw = raw.split(",", 1)[1]
        # Remove whitespace and line breaks that some platforms inject
        raw = raw.strip().replace("\n", "").replace("\r", "").replace(" ", "")
        # Fix padding
        raw += "=" * (-len(raw) % 4)
        file_bytes = base64.b64decode(raw)
        logger.info("Base64 decoded: %d bytes, magic=%r", len(file_bytes), file_bytes[:8])
    except Exception:
        raise HTTPException(status_code=400, detail="'file' is not valid base64 content.")

    filename = body.filename or "document.pdf"
    ext = Path(filename).suffix.lower()

    logger.info("Base64 upload: %s (%d bytes) | %s → %s", filename, len(file_bytes), source_lang, target_lang)

    token = token_manager.get_token()
    fn = _translate_fn(model_id, token, watsonx_url, wx_project_id, source_lang, target_lang)

    page_details: list[TranslationPageDetail] = []
    output_ext = ext if ext in SUPPORTED_EXTENSIONS else ".pdf"

    if ext == ".docx":
        output_bytes = translate_docx(file_bytes, fn)
    elif ext in (".xlsx", ".xls"):
        output_bytes = translate_xlsx(file_bytes, fn)
        output_ext = ".xlsx"
    elif ext in (".pptx", ".ppt"):
        output_bytes = translate_pptx(file_bytes, fn)
        output_ext = ".pptx"
    else:
        pages = extract_pages(file_bytes, filename, source_lang)
        if not pages:
            raise HTTPException(status_code=422, detail="No extractable text found.")

        translated_pages: list[str] = []
        for i, page_text in enumerate(pages):
            translated = translate_page(page_text, model_id, token, watsonx_url, wx_project_id, source_lang, target_lang)
            translated_pages.append(translated)
            page_details.append(TranslationPageDetail(page=i + 1, source_chars=len(page_text), translated_chars=len(translated)))

        output_ext = ".pdf"
        tmp_pdf = os.path.join(tempfile.gettempdir(), _output_filename(filename, ".pdf"))
        build_translated_pdf(translated_pages, tmp_pdf, source_lang, target_lang)
        with open(tmp_pdf, "rb") as f:
            output_bytes = f.read()

    download_url = save_and_finalize(output_bytes, filename, output_ext, str(request.base_url))

    return TranslateResponse(
        message=f"Translation complete ({source_lang} → {target_lang})",
        source_lang=source_lang,
        target_lang=target_lang,
        pages_translated=len(page_details) or 1,
        model_used=model_id,
        region=watsonx_url,
        download_url=download_url,
        pages=page_details,
    )


@app.post(
    "/api/v1/translate/from-source",
    response_model=TranslateResponse,
    tags=["Translation"],
    summary="Translate a document from a file path, URL, or bucket",
    description="Fetch a document from a server path, HTTP URL, or S3-compatible bucket and translate it.",
)
async def translate_from_source(request: Request, body: TranslateFromSourceRequest):
    model_id      = body.model_id or ModelID.GRANITE_3_8B_INSTRUCT.value
    watsonx_url   = body.region or DEFAULT_WATSONX_URL
    wx_project_id = body.project_id or WATSONX_PROJECT_ID
    source_lang   = body.source_lang or "Japanese"
    target_lang   = body.target_lang or "English"

    if not wx_project_id:
        raise HTTPException(status_code=400, detail="Project ID required.")

    file_bytes, filename = load_bytes_from_source(body.source)
    ext = Path(filename).suffix.lower()

    logger.info("From-source: %s (%d bytes) | %s → %s", filename, len(file_bytes), source_lang, target_lang)

    token = token_manager.get_token()
    fn = _translate_fn(model_id, token, watsonx_url, wx_project_id, source_lang, target_lang)

    page_details: list[TranslationPageDetail] = []
    output_ext = ext if ext in SUPPORTED_EXTENSIONS else ".pdf"

    if ext == ".docx":
        output_bytes = translate_docx(file_bytes, fn)
    elif ext in (".xlsx", ".xls"):
        output_bytes = translate_xlsx(file_bytes, fn)
        output_ext = ".xlsx"
    elif ext in (".pptx", ".ppt"):
        output_bytes = translate_pptx(file_bytes, fn)
        output_ext = ".pptx"
    else:
        pages = extract_pages(file_bytes, filename, source_lang)
        if not pages:
            raise HTTPException(status_code=422, detail="No extractable text found.")

        translated_pages: list[str] = []
        for i, page_text in enumerate(pages):
            logger.info("Translating page %d/%d...", i + 1, len(pages))
            translated = translate_page(page_text, model_id, token, watsonx_url, wx_project_id, source_lang, target_lang)
            translated_pages.append(translated)
            page_details.append(TranslationPageDetail(page=i + 1, source_chars=len(page_text), translated_chars=len(translated)))

        output_ext = ".pdf"
        tmp_pdf = os.path.join(tempfile.gettempdir(), _output_filename(filename, ".pdf"))
        build_translated_pdf(translated_pages, tmp_pdf, source_lang, target_lang)
        with open(tmp_pdf, "rb") as f:
            output_bytes = f.read()

    download_url = save_and_finalize(output_bytes, filename, output_ext, str(request.base_url))

    return TranslateResponse(
        message=f"Translation complete ({source_lang} → {target_lang})",
        source_lang=source_lang,
        target_lang=target_lang,
        pages_translated=len(page_details) or 1,
        model_used=model_id,
        region=watsonx_url,
        download_url=download_url,
        pages=page_details,
    )


@app.post(
    "/api/v1/translate/text",
    response_model=TranslateTextResponse,
    tags=["Translation"],
    summary="Translate raw text (any language pair)",
    description="Send raw text and receive the translation. Supports any source/target language.",
)
async def translate_text_endpoint(request: Request, body: TranslateTextRequest):
    model_id      = body.model_id or ModelID.GRANITE_3_8B_INSTRUCT.value
    watsonx_url   = body.region or DEFAULT_WATSONX_URL
    wx_project_id = WATSONX_PROJECT_ID
    source_lang   = body.source_lang or "auto"
    target_lang   = body.target_lang or "English"

    if not wx_project_id:
        raise HTTPException(status_code=400, detail="WATSONX_PROJECT_ID env var not set.")
    if not body.text or not body.text.strip():
        raise HTTPException(status_code=400, detail="'text' field is empty.")

    token = token_manager.get_token()
    translated = translate_page(body.text, model_id, token, watsonx_url, wx_project_id, source_lang, target_lang)

    tmp_pdf = os.path.join(tempfile.gettempdir(), _output_filename(body.filename or "document", ".pdf"))
    build_translated_pdf([translated], tmp_pdf, source_lang, target_lang)
    with open(tmp_pdf, "rb") as f:
        pdf_bytes = f.read()
    download_url = save_and_finalize(pdf_bytes, body.filename or "document", ".pdf", str(request.base_url))

    return TranslateTextResponse(
        translated_text=translated,
        source_lang=source_lang,
        target_lang=target_lang,
        model_used=model_id,
        source_chars=len(body.text),
        translated_chars=len(translated),
        download_url=download_url,
    )


# ── Download endpoint ────────────────────────────────────────────────

@app.get(
    "/api/v1/download/{filename}",
    tags=["Download"],
    summary="Download a translated document",
    responses={
        200: {"description": "Translated document (PDF, DOCX, XLSX, PPTX, …)"},
        404: {"description": "File not found"},
    },
)
async def download_translated_file(filename: str):
    safe_name = re.sub(r"[^a-zA-Z0-9._\-]", "", filename)
    path = os.path.join(tempfile.gettempdir(), safe_name)

    if not os.path.exists(path):
        raise HTTPException(status_code=404, detail="File not found or expired.")

    ext = Path(safe_name).suffix.lower()
    media_type = MIME_TYPES.get(ext, "application/octet-stream")

    return FileResponse(
        path,
        media_type=media_type,
        filename=safe_name,
        headers={"Content-Disposition": f"attachment; filename={safe_name}"},
    )


# ── Entry Point ──────────────────────────────────────────────────────

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000, log_level="info")
